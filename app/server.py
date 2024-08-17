from flask import Flask, request, jsonify
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import BedrockEmbeddings
from langchain.vectorstores import FAISS
from langchain.indexes import VectorstoreIndexCreator
from langchain.llms import Bedrock
from pptx import Presentation
import os

app = Flask(__name__)

# Default route
@app.route('/')
def home():
    return "Welcome to the RAG API server. Please use the appropriate endpoints."

# Favicon route (optional)
@app.route('/favicon.ico')
def favicon():
    return '', 204  # No content

# Custom loader for PPTX files
def load_pptx(file_path):
    print(file_path)
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

# Function to load and split documents
def load_document(file_path):
    print(f"Received file path: {file_path}")
    if file_path.endswith('.pdf'):
        loader = PyPDFLoader(file_path)
        return loader.load_and_split()
    elif file_path.endswith('.pptx'):
        pptx_text = load_pptx(file_path)
        return [{"text": pptx_text}]  # Return as a single chunk
    else:
        raise ValueError("Unsupported file type. Please upload a PDF or PPTX file.")

# Function to create index
def get_index(file_path):
    data_load = load_document(file_path)

    data_split = RecursiveCharacterTextSplitter(
        separators=["\n\n", "\n", " ", ""],
        chunk_size=100,
        chunk_overlap=10
    )

    data_embeddings = BedrockEmbeddings(
        credential_profile_name='default',
        model_id='amazon.titan-embed-text-v1'
    )

    data_index = VectorstoreIndexCreator(
        text_splitter=data_split,
        embedding=data_embeddings,
        vectorstore_cls=FAISS
    )

    db_index = data_index.from_loaders([data_load])
    return db_index

def chat_llm():
    llm = Bedrock(
        credential_profile_name='default',
        model_id='anthropic.claude-v2',
        model_kwargs={
            "max_tokens_to_sample": 300,
            "temperature": 0.1,
            "top_p": 0.9,
        }
    )
    return llm

def rag_response(index, question):
    rag_llm = chat_llm()
    chat_rag_query = index.query(question=question, llm=rag_llm)
    return chat_rag_query

@app.route('/index', methods=['POST'])
def index():
    file_path = request.json['filePath']
    try:
        print(f"Received file path: {file_path}")
        if not os.path.exists(file_path):
            print(f"File does not exist at the provided path: {file_path}")
            raise FileNotFoundError(f"File not found: {file_path}")
        
        index_data = get_index(file_path)
        print("Index created successfully.")
        return jsonify({'index': index_data})
    except Exception as e:
        print(f"Error while creating index: {str(e)}")
        return jsonify({'error': f'Failed to create index: {str(e)}'}), 500

@app.route('/respond', methods=['POST'])
def respond():
    index = request.json['index']
    question = request.json['question']
    response_data = rag_response(index, question)
    return jsonify({'response': response_data})

if __name__ == '__main__':
    app.run(debug=True)